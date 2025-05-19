import google.generativeai as genai
import pandas as pd
import os
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from datetime import datetime

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

# Function to generate slide text using Gemini
def generate_slide_text_from_data(df: pd.DataFrame, query_goal: str) -> str:
    data_description = df.to_string(index=False)
    prompt = f"""
You are a professional data analyst tasked with preparing a boardroom-level presentation based on the data below.

### Dataset:
{data_description}

### Goal:
{query_goal}

### Output Format:
Slide: <Slide Title>
• Bullet point 1
• Bullet point 2
...

### Instructions:
- Use concise, data-driven bullet points.
- Avoid unnecessary fluff or repetition.
- Include slides like Introduction, Key Insights, Visual Summary, and Conclusion.
- Output must be plain text only as described.
"""
    model = genai.GenerativeModel("gemini-1.5-flash")
    response = model.generate_content(prompt)
    return response.text.strip()


# Function to convert slide text into a styled presentation
def generate_styled_ppt(slide_text: str, output_path: str, graphs: list, author: str = "Suji19", title: str = "Generated Presentation"):
    prs = Presentation()

    font_name = "Segoe UI"
    bullet_font_size = Pt(20)
    title_font_size = Pt(36)
    subtitle_font_size = Pt(20)

    text_color = RGBColor(30, 30, 30)
    accent_color = RGBColor(0, 102, 204)

    # Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = title_slide.shapes.title
    subtitle_shape = title_slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = f"Prepared by {author}\n{datetime.today().strftime('%B %d, %Y')}"

    for paragraph in subtitle_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = subtitle_font_size
            run.font.color.rgb = text_color

    title_shape.text_frame.paragraphs[0].runs[0].font.size = title_font_size
    title_shape.text_frame.paragraphs[0].runs[0].font.name = font_name
    title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = accent_color

    slide_layout = prs.slide_layouts[5]  # Title + content

    def add_slide(title, content_lines):
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        title_shape.text = title
        title_run = title_shape.text_frame.paragraphs[0].runs[0]
        title_run.font.size = title_font_size
        title_run.font.name = font_name
        title_run.font.color.rgb = accent_color

        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        textbox = shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        for line in content_lines:
            if line.strip():
                p = text_frame.add_paragraph()
                p.text = line.strip("• ").strip("- ")
                p.level = 0
                p.font.size = bullet_font_size
                p.font.name = font_name
                p.font.color.rgb = text_color

    # Parse and add slides
    current_title = None
    current_content = []

    for line in slide_text.splitlines():
        line = line.strip()
        if not line:
            continue
        if line.lower().startswith("slide:"):
            if current_title and current_content:
                add_slide(current_title, current_content)
            current_title = line[6:].strip()
            current_content = []
        else:
            current_content.append(line)

    if current_title and current_content:
        add_slide(current_title, current_content)
    
    print("graphs in gen_ppt : ", graphs)

    # Add Graph Slides if any
    if graphs:
        for idx, graph_path in enumerate(graphs, 1):
            try:
                graph_slide = prs.slides.add_slide(prs.slide_layouts[5])
                title_shape = graph_slide.shapes.title
                title_shape.text = f"Graph {idx}"

                left = Inches(1)
                top = Inches(1.5)
                width = Inches(8)
                height = Inches(4.5)

                if os.path.exists(graph_path):
                    graph_slide.shapes.add_picture(graph_path, left, top, width=width, height=height)
                else:
                    textbox = graph_slide.shapes.add_textbox(left, top, width, height)
                    tf = textbox.text_frame
                    tf.text = f"⚠️ Could not find image at:\n{graph_path}"
            except Exception as e:
                print(f"🚨 Error embedding graph {graph_path}: {e}")

    prs.save(output_path)
    print(f"✅ Presentation saved as '{output_path}' with {len(prs.slides)} slides.")


# Master Function to call from outside
def generate_presentation(df: pd.DataFrame, query_goal: str, output_path: str,graphs: list, author: str = "Sujnan", title: str = "Presentation"):
    slide_text = generate_slide_text_from_data(df, query_goal)
    print("Graph in main func : ", graphs)
    generate_styled_ppt(slide_text, output_path, graphs , author, title)
